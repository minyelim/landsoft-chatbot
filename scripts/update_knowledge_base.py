"""
회사 규정 PDF 또는 온보딩 PPT 파일이 새로 업데이트되었을 때,
knowledge_base 폴더의 텍스트 파일을 다시 만들어주는 스크립트입니다.

사용 예시 (프로젝트 폴더에서 실행):
    python scripts/update_knowledge_base.py --pdf "회사규정_2025_01.pdf" --out 01_회사규정
    python scripts/update_knowledge_base.py --pptx "신규입사자_교육자료.pptx" --out 02_온보딩가이드

--out 뒤에는 파일명(확장자 제외)만 적어주세요. knowledge_base 폴더 안에
"해당이름.txt" 로 저장됩니다. 기존 파일과 같은 이름을 쓰면 덮어씁니다.
"""

import argparse
from pathlib import Path

KB_DIR = Path(__file__).parent.parent / "knowledge_base"


def extract_pdf(pdf_path: str) -> str:
    import pdfplumber

    lines = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            lines.append(f"\n===PAGE {i}===\n{text}")
    return "\n".join(lines)


def extract_pptx(pptx_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n===SLIDE {i}===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(" | ".join(cells))
    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--pdf", help="새 규정 PDF 파일 경로")
    parser.add_argument("--pptx", help="새 온보딩 PPTX 파일 경로")
    parser.add_argument("--out", required=True, help="저장할 파일명 (확장자 제외)")
    args = parser.parse_args()

    if not args.pdf and not args.pptx:
        raise SystemExit("--pdf 또는 --pptx 중 하나는 반드시 지정해야 합니다.")

    if args.pdf:
        content = extract_pdf(args.pdf)
    else:
        content = extract_pptx(args.pptx)

    KB_DIR.mkdir(exist_ok=True)
    out_path = KB_DIR / f"{args.out}.txt"
    out_path.write_text(content, encoding="utf-8")
    print(f"저장 완료: {out_path}")


if __name__ == "__main__":
    main()
